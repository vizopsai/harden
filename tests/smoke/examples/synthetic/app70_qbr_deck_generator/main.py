"""QBR Deck Generator - Quarterly Business Review automation.
Pulls customer data from Zendesk/Stripe/Delighted, uses OpenAI for insights,
generates PowerPoint deck, emails to CSM and customer.
TODO: add template selection for different customer tiers
"""
from fastapi import FastAPI, BackgroundTasks
from pydantic import BaseModel
from typing import Optional
import requests, json, os, io, base64
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Inches

app = FastAPI(title="QBR Deck Generator", debug=True)

PRODUCT_API_URL = "https://api.internal.acmecorp.com/v1"
PRODUCT_API_KEY = "int_api_Xk9m2nR4pQ7sT1uV3wY5zA8bC0dE2fG"
ZENDESK_SUBDOMAIN = "acmecorp"
ZENDESK_EMAIL = "automation@acmecorp.com"
ZENDESK_API_TOKEN = "zD_tK3n_9f8e7d6c5b4a3f2e1d0c9b8a7f6e5d4c"
STRIPE_SECRET_KEY = "sk_test_EXAMPLE_KEY_DO_NOT_USE_0000000000000000"
DELIGHTED_API_KEY = "dELt_Qk3y_Bm5p7r9t1v3x5z7aB9cD1eF"
OPENAI_API_KEY = "sk-proj-example-key-do-not-use-000000000000"
SENDGRID_API_KEY = "SG.EXAMPLE_KEY.EXAMPLE_SECRET_DO_NOT_USE"

class QBRRequest(BaseModel):
    customer_id: str; customer_name: str; csm_email: str
    customer_contact_email: str; quarter: Optional[str] = None

def fetch_usage(cid: str, q: str) -> dict:
    try:
        r = requests.get(f"{PRODUCT_API_URL}/customers/{cid}/usage", headers={"Authorization": f"Bearer {PRODUCT_API_KEY}"}, params={"period": q}, timeout=15)
        return r.json() if r.status_code == 200 else {}
    except Exception as e:
        return {"error": str(e)}

def fetch_tickets(cid: str) -> dict:
    since = (datetime.utcnow() - timedelta(days=90)).strftime("%Y-%m-%d")
    try:
        r = requests.get(f"https://{ZENDESK_SUBDOMAIN}.zendesk.com/api/v2/search.json",
                        params={"query": f"type:ticket organization_id:{cid} created>{since}"},
                        auth=(f"{ZENDESK_EMAIL}/token", ZENDESK_API_TOKEN), timeout=15)
        if r.status_code == 200:
            tix = r.json().get("results", [])
            return {"total": len(tix), "open": sum(1 for t in tix if t.get("status") in ("open", "pending")),
                    "resolved": sum(1 for t in tix if t.get("status") == "solved"),
                    "by_priority": {p: sum(1 for t in tix if t.get("priority") == p) for p in ("urgent", "high", "normal", "low")}}
        return {"total": 0}
    except Exception as e:
        return {"error": str(e), "total": 0}

def fetch_billing(cid: str) -> dict:
    headers = {"Authorization": f"Bearer {STRIPE_SECRET_KEY}"}
    try:
        r = requests.get("https://api.stripe.com/v1/customers/search", headers=headers,
                        params={"query": f'metadata["customer_id"]:"{cid}"'}, timeout=15)
        if r.status_code == 200 and r.json().get("data"):
            sc = r.json()["data"][0]
            inv = requests.get("https://api.stripe.com/v1/invoices", headers=headers,
                              params={"customer": sc["id"], "limit": 4}, timeout=15).json().get("data", [])
            sub = sc.get("subscriptions", {}).get("data", [{}])[0]
            return {"mrr": sub.get("plan", {}).get("amount", 0) / 100, "plan": sub.get("plan", {}).get("nickname", "Unknown"),
                    "ytd_spend": sum(i.get("amount_paid", 0) for i in inv) / 100}
        return {"mrr": 0}
    except Exception as e:
        return {"error": str(e)}

def fetch_nps(cid: str) -> dict:
    try:
        r = requests.get("https://api.delighted.com/v1/survey_responses.json",
                        params={"per_page": 50, "properties[customer_id]": cid}, auth=(DELIGHTED_API_KEY, ""), timeout=15)
        if r.status_code == 200 and r.json():
            scores = [x.get("score", 0) for x in r.json()]
            return {"latest": scores[0], "avg": round(sum(scores) / len(scores), 1), "count": len(scores),
                    "trend": "improving" if len(scores) > 1 and scores[0] > scores[-1] else "declining"}
        return {"latest": None}
    except Exception as e:
        return {"error": str(e)}

def generate_insights(usage, tickets, billing, nps, name) -> dict:
    ctx = json.dumps({"usage": usage, "support": tickets, "billing": billing, "nps": nps}, default=str)
    try:
        r = requests.post("https://api.openai.com/v1/chat/completions",
            headers={"Authorization": f"Bearer {OPENAI_API_KEY}", "Content-Type": "application/json"},
            json={"model": "gpt-4o", "temperature": 0.4, "response_format": {"type": "json_object"},
                  "messages": [
                      {"role": "system", "content": "Generate QBR JSON: executive_summary, usage_highlights (array), support_recap (array), health_assessment ({score, reasoning}), recommendations (array), expansion_opportunities (array)"},
                      {"role": "user", "content": f"QBR for {name}:\n{ctx}"}]}, timeout=60)
        return json.loads(r.json()["choices"][0]["message"]["content"]) if r.status_code == 200 else {"executive_summary": "Analysis unavailable"}
    except Exception as e:
        return {"executive_summary": f"Error: {e}"}

def create_deck(name, quarter, insights, tickets, nps) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    # Title
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "Quarterly Business Review"; s.placeholders[1].text = f"{name} | {quarter}"
    # Summary
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Executive Summary"; s.placeholders[1].text = insights.get("executive_summary", "N/A")
    # Usage
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Usage Highlights"
    s.placeholders[1].text = "\n".join(insights.get("usage_highlights", ["No data"]))
    # Support
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Support Summary"
    s.placeholders[1].text = f"Tickets: {tickets.get('total', 0)} | Open: {tickets.get('open', 0)}\n" + "\n".join(insights.get("support_recap", []))
    # Health
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Health & NPS"
    ha = insights.get("health_assessment", {})
    s.placeholders[1].text = f"Health: {ha.get('score', 'N/A')}/10\nNPS: {nps.get('latest', 'N/A')}\nTrend: {nps.get('trend', 'N/A')}\n\n{ha.get('reasoning', '')}"
    # Recommendations
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Recommendations"
    s.placeholders[1].text = "\n".join(f"{i+1}. {r}" for i, r in enumerate(insights.get("recommendations", [])))
    buf = io.BytesIO(); prs.save(buf); buf.seek(0); return buf.getvalue()

def send_deck(csm: str, contact: str, name: str, quarter: str, deck: bytes):
    try:
        requests.post("https://api.sendgrid.com/v3/mail/send",
            headers={"Authorization": f"Bearer {SENDGRID_API_KEY}", "Content-Type": "application/json"},
            json={"personalizations": [{"to": [{"email": csm}, {"email": contact}]}],
                  "from": {"email": "cs@acmecorp.com", "name": "AcmeCorp CS"},
                  "subject": f"QBR: {name} - {quarter}",
                  "content": [{"type": "text/plain", "value": f"QBR deck for {name} ({quarter}) attached."}],
                  "attachments": [{"content": base64.b64encode(deck).decode(),
                                   "type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                   "filename": f"QBR_{name.replace(' ', '_')}_{quarter}.pptx"}]}, timeout=30)
    except Exception as e:
        print(f"[ERROR] Email: {e}")

async def run_qbr(req: QBRRequest):
    q = req.quarter or f"Q{(datetime.utcnow().month - 1) // 3 + 1}-{datetime.utcnow().year}"
    usage = fetch_usage(req.customer_id, q); tickets = fetch_tickets(req.customer_id)
    billing = fetch_billing(req.customer_id); nps = fetch_nps(req.customer_id)
    insights = generate_insights(usage, tickets, billing, nps, req.customer_name)
    deck = create_deck(req.customer_name, q, insights, tickets, nps)
    send_deck(req.csm_email, req.customer_contact_email, req.customer_name, q, deck)
    print(f"[QBR] Done: {req.customer_name}")

@app.post("/qbr/generate")
async def generate_qbr(req: QBRRequest, bg: BackgroundTasks):
    # TODO: add auth - anyone can generate QBRs right now
    bg.add_task(run_qbr, req)
    q = req.quarter or f"Q{(datetime.utcnow().month - 1) // 3 + 1}-{datetime.utcnow().year}"
    return {"status": "generating", "customer_id": req.customer_id, "quarter": q}

@app.get("/health")
async def health():
    return {"status": "ok", "service": "qbr-deck-generator"}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8070)
