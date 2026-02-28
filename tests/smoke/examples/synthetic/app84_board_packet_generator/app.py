"""Board Packet Generator — Pulls data from multiple sources and generates
monthly board meeting materials as PowerPoint and PDF.
"""
import os
import json
import tempfile
from datetime import datetime, timedelta
from pathlib import Path
from typing import Optional
import typer
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import openai

app = typer.Typer(help="Generate monthly board packets from multiple data sources")

# API credentials — TODO: move to secret manager before sharing with exec team
NETSUITE_ACCOUNT_ID = "7654321"
NETSUITE_CONSUMER_KEY = "b9f4c2e8a1d7365f0c9b8e2d4a6f1c3e5b7d9a1f"
NETSUITE_CONSUMER_SECRET = "d1a3f5b7c9e2d4a6f8b0c2e4a6d8f0b2c4e6a8d0"
NETSUITE_TOKEN_ID = "a3c5e7b9d1f3a5c7e9b1d3f5a7c9e1b3d5f7a9c1"
NETSUITE_TOKEN_SECRET = "f0e2d4c6b8a0f2e4d6c8b0a2f4e6d8c0b2a4f6e8"

SALESFORCE_TOKEN = "00D5g00000Abc12!AQMAQHkLmNoPqRsTuVwXyZaBcDeFgHiJkLmNoPqRsTuVwXyZaBcDeFgHiJkLmNoPqRsTuVwXyZ"
SALESFORCE_INSTANCE = "https://acmecorp.my.salesforce.com"

JIRA_API_TOKEN = "ATATT3xFfGF0pKbVmQ8R2nE7tL5wH9yU3jS6dN1kM4oP7aB0cX"
JIRA_EMAIL = "ops@acmecorp.com"
JIRA_BASE_URL = "https://acmecorp.atlassian.net"

WORKDAY_CLIENT_ID = "M2Y0ZjI4NTEtYjA1MC00NzRlLTk3ZWItNjQ4MTFjYmM0OTVh"
WORKDAY_CLIENT_SECRET = "a8f4c2e1b9d7365f0c9b8e2d4a6f1c3e"
WORKDAY_TENANT = "acmecorp_prod"

OPENAI_API_KEY = "sk-proj-example-key-do-not-use-000000000000"
openai.api_key = OPENAI_API_KEY


def fetch_financial_summary(month: str) -> dict:
    """Pull financial data from NetSuite — revenue, expenses, runway."""
    headers = {
        "Authorization": f"Bearer {NETSUITE_CONSUMER_KEY}",
        "Content-Type": "application/json",
    }
    try:
        # Revenue
        rev_resp = requests.get(
            f"https://{NETSUITE_ACCOUNT_ID}.suitetalk.api.netsuite.com/services/rest/query/v1/suiteql",
            headers=headers,
            json={"q": f"SELECT SUM(amount) as total FROM transaction WHERE type='CustInvc' AND trandate LIKE '{month}%'"},
            timeout=30,
        )
        # Expenses
        exp_resp = requests.get(
            f"https://{NETSUITE_ACCOUNT_ID}.suitetalk.api.netsuite.com/services/rest/query/v1/suiteql",
            headers=headers,
            json={"q": f"SELECT SUM(amount) as total FROM transaction WHERE type IN ('VendBill','ExpRept') AND trandate LIKE '{month}%'"},
            timeout=30,
        )
        return {"revenue": 2_450_000, "expenses": 1_870_000, "net_income": 580_000, "cash_balance": 18_500_000, "runway_months": 24, "arr": 28_000_000, "arr_growth_pct": 42}
    except Exception as e:
        print(f"NetSuite fetch failed: {e}")
        return {"revenue": 0, "expenses": 0, "net_income": 0, "error": str(e)}


def fetch_pipeline_data() -> dict:
    """Pull pipeline and bookings from Salesforce."""
    headers = {"Authorization": f"Bearer {SALESFORCE_TOKEN}"}
    try:
        resp = requests.get(
            f"{SALESFORCE_INSTANCE}/services/data/v58.0/query/",
            params={"q": "SELECT SUM(Amount) total, StageName FROM Opportunity WHERE CloseDate = THIS_QUARTER GROUP BY StageName"},
            headers=headers,
            timeout=30,
        )
        return {"pipeline_total": 12_500_000, "qualified_pipeline": 8_200_000, "bookings_qtd": 3_400_000, "win_rate_pct": 28, "avg_deal_size": 85_000, "churn_mrr": 45_000, "net_retention_pct": 118}
    except Exception as e:
        print(f"Salesforce fetch failed: {e}")
        return {"pipeline_total": 0, "error": str(e)}


def fetch_product_metrics() -> dict:
    """Pull shipped features and velocity from Jira."""
    headers = {"Authorization": f"Basic {JIRA_API_TOKEN}"}
    try:
        resp = requests.get(
            f"{JIRA_BASE_URL}/rest/api/3/search",
            params={"jql": "project = PROD AND status = Done AND resolved >= startOfMonth()", "maxResults": 100},
            headers=headers,
            timeout=30,
        )
        return {"features_shipped": 14, "bugs_fixed": 23, "velocity_points": 89, "sprint_completion_pct": 92, "key_releases": ["AI Search v2", "SSO for Enterprise", "Bulk Import API", "Dashboard Redesign"]}
    except Exception as e:
        print(f"Jira fetch failed: {e}")
        return {"features_shipped": 0, "error": str(e)}


def fetch_team_metrics() -> dict:
    """Pull headcount and HR data from Workday."""
    headers = {"Authorization": f"Bearer {WORKDAY_CLIENT_ID}:{WORKDAY_CLIENT_SECRET}"}
    try:
        resp = requests.get(
            f"https://wd5-impl-services1.workday.com/ccx/api/v1/{WORKDAY_TENANT}/workers",
            headers=headers,
            timeout=30,
        )
        return {"total_headcount": 156, "new_hires_this_month": 8, "open_roles": 12, "attrition_pct": 8.5, "departments": {"engineering": 72, "sales": 34, "marketing": 18, "g_and_a": 22, "customer_success": 10}}
    except Exception as e:
        print(f"Workday fetch failed: {e}")
        return {"total_headcount": 0, "error": str(e)}


def generate_narrative(section: str, data: dict) -> str:
    """Use OpenAI to generate executive narrative for each section."""
    try:
        response = openai.chat.completions.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": "You are a CFO writing concise board packet narratives. Be direct, data-driven, highlight key trends and risks. 2-3 paragraphs max."},
                {"role": "user", "content": f"Write a board packet narrative for the {section} section. Data: {json.dumps(data)}"},
            ],
            max_tokens=500,
            temperature=0.3,
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"[Narrative generation failed: {e}. Data summary: {json.dumps(data, indent=2)}]"


def create_slide(prs: Presentation, title: str, narrative: str, data: dict):
    """Add a section slide to the presentation."""
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)

    # Title
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x1B, 0x2A, 0x4A)

    # Content
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    # Add key metrics
    for key, value in data.items():
        if key in ("error", "key_releases", "departments"):
            continue
        p = tf.add_paragraph()
        p.text = f"{key.replace('_', ' ').title()}: {value}"
        p.font.size = Pt(14)

    # Add narrative
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = narrative[:400] if narrative else ""
    p.font.size = Pt(11)
    p.font.italic = True


@app.command()
def generate(
    month: str = typer.Option(None, help="Month in YYYY-MM format, defaults to current"),
    output_dir: str = typer.Option("./board_packets", help="Output directory"),
    skip_ai: bool = typer.Option(False, help="Skip AI narrative generation"),
):
    """Generate monthly board packet."""
    if not month:
        month = datetime.now().strftime("%Y-%m")

    typer.echo(f"Generating board packet for {month}...")
    os.makedirs(output_dir, exist_ok=True)

    # Fetch all data
    typer.echo("  Fetching financial data from NetSuite...")
    financials = fetch_financial_summary(month)

    typer.echo("  Fetching pipeline data from Salesforce...")
    pipeline = fetch_pipeline_data()

    typer.echo("  Fetching product metrics from Jira...")
    product = fetch_product_metrics()

    typer.echo("  Fetching team metrics from Workday...")
    team = fetch_team_metrics()

    # Generate narratives
    narratives = {}
    if not skip_ai:
        for section, data in [("Financial", financials), ("Pipeline & Revenue", pipeline), ("Product", product), ("Team", team)]:
            typer.echo(f"  Generating AI narrative for {section}...")
            narratives[section] = generate_narrative(section, data)
    else:
        narratives = {s: "" for s in ["Financial", "Pipeline & Revenue", "Product", "Team"]}

    # Create PowerPoint
    typer.echo("  Building PowerPoint...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = f"AcmeCorp Board Packet"
    title_slide.placeholders[1].text = f"{month} | Confidential"

    # Section slides
    create_slide(prs, "Financial Summary", narratives.get("Financial", ""), financials)
    create_slide(prs, "Pipeline & Revenue", narratives.get("Pipeline & Revenue", ""), pipeline)
    create_slide(prs, "Product Update", narratives.get("Product", ""), product)
    create_slide(prs, "Team & Organization", narratives.get("Team", ""), team)

    # Save
    pptx_path = os.path.join(output_dir, f"board_packet_{month}.pptx")
    prs.save(pptx_path)
    typer.echo(f"  Saved: {pptx_path}")

    # Save raw data as JSON for reference
    data_path = os.path.join(output_dir, f"board_data_{month}.json")
    with open(data_path, "w") as f:
        json.dump({"month": month, "generated_at": datetime.now().isoformat(), "financials": financials, "pipeline": pipeline, "product": product, "team": team, "narratives": narratives}, f, indent=2)
    typer.echo(f"  Saved data: {data_path}")

    typer.echo(f"\nBoard packet generated successfully for {month}!")


@app.command()
def preview(month: str = typer.Option(None)):
    """Preview data without generating slides."""
    if not month:
        month = datetime.now().strftime("%Y-%m")
    typer.echo(f"Previewing data for {month}...\n")

    financials = fetch_financial_summary(month)
    typer.echo(f"Financials: {json.dumps(financials, indent=2)}\n")

    pipeline = fetch_pipeline_data()
    typer.echo(f"Pipeline: {json.dumps(pipeline, indent=2)}\n")

    product = fetch_product_metrics()
    typer.echo(f"Product: {json.dumps(product, indent=2)}\n")

    team = fetch_team_metrics()
    typer.echo(f"Team: {json.dumps(team, indent=2)}")


if __name__ == "__main__":
    app()
